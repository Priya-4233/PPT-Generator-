from llm_utils import generate_titles, generate_bullets
from image_utils import get_image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

import requests
from io import BytesIO
import time


# ---------------- TEXT CLEANER ----------------
def clean_bullet(text, max_words=12):
    words = str(text).split()
    if len(words) > max_words:
        return " ".join(words[:max_words]) + "..."
    return " ".join(words)


# ---------------- PPT CREATOR ----------------
def create_ppt(topic: str, filename=None):

    print("\n🧠 Generating slide titles...")
    titles = generate_titles(topic)

    print("\n🧭 Plan:")
    for i, t in enumerate(titles, 1):
        print(f"{i}. {t}")

    if filename is None:
        safe_topic = topic[:30].replace(" ", "_").replace(":", "")
        filename = f"PPT_{safe_topic}_{int(time.time())}.pptx"

    prs = Presentation()

    for title in titles:
        print(f"\n➕ Creating slide: {title}")

        bullets = generate_bullets(topic, title)

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # TITLE
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)
        )

        title_tf = title_box.text_frame
        title_tf.text = title

        title_para = title_tf.paragraphs[0]
        title_para.font.size = Pt(30)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # BULLETS
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2), Inches(4.8), Inches(5.2)
        )

        tf = content_box.text_frame
        tf.clear()

        for i, b in enumerate(bullets):
            b = clean_bullet(b)

            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.font.size = Pt(16)
            p.level = 0

        # IMAGE
        img_url = get_image(f"{topic} {title}")

        if img_url:
            try:
                img_data = requests.get(img_url, timeout=10).content
                img_stream = BytesIO(img_data)

                slide.shapes.add_picture(
                    img_stream,
                    Inches(5.4),
                    Inches(1.2),
                    width=Inches(4.0),
                    height=Inches(4.5)
                )
            except Exception as e:
                print("⚠️ Image insert failed:", e)

    prs.save(filename)

    print(f"\n✅ PPT successfully saved: {filename}")


# ---------------- MAIN ----------------
if __name__ == "__main__":
    topic = input("Enter PPT topic: ")
    create_ppt(topic)